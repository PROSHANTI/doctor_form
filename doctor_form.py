#!/usr/bin/env python3
import datetime
import locale
import os

from tkinter import *
from tkinter import messagebox
import tkinter as tk
from tkinter import ttk
from pptx import Presentation

locale.setlocale(locale.LC_ALL, 'ru_RU.UTF-8')


def gui():
    def get_current_date():
        """Получить текущую дату и вернуть день, месяц и год."""
        now = datetime.datetime.now()
        return now.day, now.strftime("%B"), now.year

    def check_input_file(input_file_path):
        """Проверяет существование файлов в папке с программой."""
        if not os.path.exists(input_file_path):
            messagebox.showinfo(
                message=f'Файл «{input_file_path}» '
                        f'не найден в папке с программой.'
            )
            return False

    def check_output_file(output_file_path):
        if os.path.exists(output_file_path):
            messagebox.showinfo(
                message=f'Файл «{output_file_path}» успешно создан'
            )
        return True

    def replace_text_in_presentation():
        """
        Заменяет текст в файле презентации пользовательским
        вводом для имён врачей и пациентов, а также текущей даты.
        """
        input_file_path = 'Бланк Врача.pptx'
        output_file_path = 'Бланк врача на печать.pptx'
        check_input_file(input_file_path)
        prs = Presentation(input_file_path)
        day, month, year = get_current_date()

        replacements = {
            'Doctor_1': f'ВРАЧ: {doctor_1.get()}',
            'Doctor_2': f'ВРАЧ: {doctor_2.get()}',
            'Doctor_3': f'ВРАЧ: {doctor_3.get()}',
            'Doctor_4': f'ВРАЧ: {doctor_4.get()}',
            'Pacient_1': f'ПАЦИЕНТ: {pacient_1.get().upper()}',
            'Pacient_2': f'ПАЦИЕНТ: {pacient_2.get().upper()}',
            'Pacient_3': f'ПАЦИЕНТ: {pacient_3.get().upper()}',
            'Pacient_4': f'ПАЦИЕНТ: {pacient_4.get().upper()}',
            'Дата': f'«{date.get()}» {month} {year} г.'
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            run.text = run.text.replace(key, value)
        prs.save(output_file_path)
        check_output_file(output_file_path)

    # GUI settings
    window = Tk()
    window.title('Cоздание красивых карточек')
    window.geometry('300x500')

    doctors = [
        'ЩУКАРЕВА Е.А.',
        'БАШИЛОВА А.А.',
        'АЛБАСТОВА М.Т.',
        'БАСКАКОВА О.А.',
        'ПУЛЬКИНА С.В.',
    ]

    frame = Frame(window)
    frame.pack()

    doctor = ttk.Label(frame, text=f"ВРАЧ")
    doctor.pack(padx=130, pady=5)

    doctor_1 = ttk.Combobox(values=doctors)
    doctor_1.pack(padx=50, pady=5)

    doctor_2 = ttk.Combobox(values=doctors)
    doctor_2.pack(padx=50, pady=5)

    doctor_3 = ttk.Combobox(values=doctors)
    doctor_3.pack(padx=50, pady=5)

    doctor_4 = ttk.Combobox(values=doctors)
    doctor_4.pack(padx=50, pady=5)

    frame = Frame(window)
    frame.pack()

    pacient = ttk.Label(frame, text='ПАЦИЕНТ')
    pacient.pack(padx=100)

    pacient_1 = ttk.Entry(frame)
    pacient_1.pack(padx=50, pady=5)

    pacient_2 = ttk.Entry(frame)
    pacient_2.pack(padx=50, pady=5)

    pacient_3 = ttk.Entry(frame)
    pacient_3.pack(padx=50, pady=5)

    pacient_4 = ttk.Entry(frame)
    pacient_4.pack(padx=50, pady=5)

    date = ttk.Label(frame, text='ДАТА')
    date.pack(padx=130)

    date = tk.Entry(frame)
    date.insert(END, get_current_date()[0])
    date.pack(padx=50, pady=5)

    btn_create = tk.Button(
        window, text='Сделать красиво!',
        command=replace_text_in_presentation,
    )
    btn_create.pack(pady=10)

    exit_button = tk.Button(window, text="Выход", command=window.destroy)
    exit_button.pack(pady=10)
    window.mainloop()


if __name__ == '__main__':
    gui()
